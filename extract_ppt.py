from pptx import Presentation
from pptx.util import Inches
import os
import json

prs = Presentation(r'd:\EvolveProfile\WEBSITE.ppt')

output_dir = r'd:\EvolveProfile\evolvewood-website\public\images\ppt_extracted'
os.makedirs(output_dir, exist_ok=True)

slides_data = []
img_counter = 0

for slide_idx, slide in enumerate(prs.slides):
    slide_info = {
        'slide_number': slide_idx + 1,
        'texts': [],
        'images': [],
        'shapes': []
    }
    
    for shape in slide.shapes:
        # Get shape info
        shape_info = {
            'name': shape.name,
            'shape_type': str(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
        }
        
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_info['texts'].append(text)
        
        if shape.shape_type == 13:  # Picture
            img_counter += 1
            image = shape.image
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            img_filename = f'slide{slide_idx+1}_img{img_counter}.{ext}'
            img_path = os.path.join(output_dir, img_filename)
            with open(img_path, 'wb') as f:
                f.write(image.blob)
            slide_info['images'].append({
                'filename': img_filename,
                'content_type': image.content_type,
                'size': len(image.blob)
            })
            shape_info['image'] = img_filename
        
        slide_info['shapes'].append(shape_info)
    
    slides_data.append(slide_info)

# Print all extracted data
for slide in slides_data:
    print(f"\n=== SLIDE {slide['slide_number']} ===")
    print("TEXTS:")
    for t in slide['texts']:
        print(f"  - {t}")
    print("IMAGES:")
    for img in slide['images']:
        print(f"  - {img['filename']} ({img['content_type']}, {img['size']} bytes)")

print(f"\nTotal images extracted: {img_counter}")
print(f"Images saved to: {output_dir}")
